# -*- coding: utf-8 -*-
"""Generate a 3-slide schematic deck on the two core innovations of proactive scanning.
Native python-pptx shapes (no rasterized boxes): rmap fan-out, MGLRU sweep, DAMON
sampling, a memory-vs-time curve, and a 4-node DAMOS control loop.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- palette (systems / kernel: slate + teal accent + red warning) ----
DARK   = RGBColor(0x0F,0x1E,0x2E)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
INK    = RGBColor(0x1E,0x29,0x3B)
MUTE   = RGBColor(0x64,0x74,0x8B)
SLATE  = RGBColor(0x94,0xA3,0xB8)
SLATEF = RGBColor(0xF1,0xF5,0xF9)
TEAL   = RGBColor(0x0D,0x94,0x88)
TEALDK = RGBColor(0x0B,0x6E,0x66)
TEALLT = RGBColor(0x99,0xE6,0xDA)
BAND   = RGBColor(0xE6,0xFB,0xF8)
RED    = RGBColor(0xD8,0x4B,0x3C)

YAHEI = "Microsoft YaHei"
MONO  = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _font(r, t, size, color, bold, font, italic):
    r.text = t
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    for tag in ('a:latin','a:ea','a:cs'):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set('typeface', font)

def fill_tf(tf, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, space_after=3):
    tf.word_wrap = wrap
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        for rs in para:
            r = p.add_run()
            _font(r, rs['t'], rs.get('s',14), rs.get('c',INK),
                  rs.get('b',False), rs.get('f',YAHEI), rs.get('i',False))

def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, space_after=3):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    fill_tf(tb.text_frame, paras, align, anchor, wrap, space_after)
    return tb

def line1(t, s, c, b=False, f=YAHEI, i=False):
    return [[{'t':t,'s':s,'c':c,'b':b,'f':f,'i':i}]]

def box(slide, shape, x, y, w, h, fill=None, ln=None, lw=1.0, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if ln is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = ln; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def connect(slide, x1, y1, x2, y2, color=TEAL, w=1.5, dash=None, arrow_end=True, arrow_begin=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    lne = cn.line._get_or_add_ln()
    if dash:
        lne.append(lne.makeelement(qn('a:prstDash'), {'val': dash}))
    if arrow_begin:
        lne.append(lne.makeelement(qn('a:headEnd'), {'type':'triangle','w':'med','len':'med'}))
    if arrow_end:
        lne.append(lne.makeelement(qn('a:tailEnd'), {'type':'triangle','w':'med','len':'med'}))
    cn.shadow.inherit = False
    return cn

# ============================ SLIDE 1 — title ============================
s = prs.slides.add_slide(BLANK); bg(s, DARK)
text(s, 0.9, 2.25, 11.6, 1.3, line1("终端「主动扫描」的两大核心创新", 40, WHITE, True), anchor=MSO_ANCHOR.TOP)
text(s, 0.92, 3.7, 11.4, 0.6, line1("机制:之前 → 之后　·　以及变化带来的收益", 18, TEALLT))
# preview chips
for cy, num, lbl in [(4.95,"1","瞄准镜变便宜、与规模解耦（扫描器）"),
                     (5.78,"2","触发与控制:被动 → 主动闭环")]:
    o = box(s, MSO_SHAPE.OVAL, 0.95, cy, 0.52, 0.52, fill=TEAL)
    fill_tf(o.text_frame, line1(num, 20, WHITE, True), PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    text(s, 1.65, cy+0.04, 11, 0.5, line1(lbl, 18, WHITE))
text(s, 0.92, 6.72, 11, 0.4, line1("基于  A16a · LRU 主动扫描  与  DAMOS 源码分析", 13, SLATE))

# ============================ SLIDE 2 — innovation 1 =====================
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
text(s, 0.5, 0.3, 12.4, 0.6, line1("创新一 · 瞄准镜变便宜、与规模解耦", 30, INK, True))
text(s, 0.5, 0.96, 12.3, 0.4, line1("扫描器升级 —— 这是「使能项」,本身还不是 proactive", 14, MUTE))
text(s, 0.55, 1.46, 5.6, 0.4, line1("之前 · 经典 LRU + rmap 反查", 16, MUTE, True))
text(s, 7.05, 1.46, 5.8, 0.4, line1("之后 · 两个低开销扫描器", 16, TEAL, True))

# --- BEFORE: rmap fan-out ---
text(s, 0.75, 1.78, 4, 0.3, line1("多个进程的页表项", 11, MUTE))
pte_ys = [2.08, 2.61, 3.14, 3.67, 4.20]
for py in pte_ys:
    r = box(s, MSO_SHAPE.RECTANGLE, 0.75, py, 1.25, 0.42, fill=WHITE, ln=SLATE, lw=1.0)
    fill_tf(r.text_frame, line1("PTE", 11, INK, False, MONO), PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
pg = box(s, MSO_SHAPE.OVAL, 4.55, 3.0, 1.3, 0.82, fill=SLATEF, ln=SLATE, lw=1.25)
fill_tf(pg.text_frame, line1("物理页", 12, INK, True), PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
for py in pte_ys:
    connect(s, 4.55, 3.41, 2.0, py+0.21, color=SLATE, w=1.1, arrow_end=True)
text(s, 2.05, 4.62, 3.6, 0.3, line1("rmap 反查:从页找回每个映射", 11, MUTE, i=True))
text(s, 0.6, 4.98, 5.5, 0.8,
     [[{'t':"✗ ",'s':12.5,'c':RED,'b':True},
       {'t':"判一页冷热要反查全部映射;成本随映射数 / 内存规模涨,大内存一次扫描达数十 ms",'s':12.5,'c':RED}]])

# --- middle upgrade arrow ---
connect(s, 6.25, 3.4, 6.98, 3.4, color=TEAL, w=3.0, arrow_end=True)
text(s, 6.18, 3.0, 0.95, 0.35, line1("升级", 14, TEAL, True), PP_ALIGN.CENTER)

# --- AFTER top: MGLRU sweep ---
text(s, 7.05, 1.86, 5.7, 0.32, line1("MGLRU · 走页表 aging", 13, TEALDK, True))
connect(s, 7.18, 2.42, 10.62, 2.42, color=TEAL, w=2.0, arrow_end=True)
text(s, 10.72, 2.26, 2.5, 0.3, line1("批量清 Accessed 位", 11, MUTE))
cx = 7.18
for _ in range(6):
    box(s, MSO_SHAPE.RECTANGLE, cx, 2.52, 0.5, 0.42, fill=WHITE, ln=TEAL, lw=1.0)
    cx += 0.58
for bx, lbl, fl, tc in [(7.2,"热代",TEALLT,INK),(8.45,"温代",RGBColor(0x5E,0xC9,0xBC),WHITE),(9.7,"冷代",TEAL,WHITE)]:
    rb = box(s, MSO_SHAPE.ROUNDED_RECTANGLE, bx, 3.06, 1.05, 0.4, fill=fl, ln=TEAL, lw=0.75, radius=0.18)
    fill_tf(rb.text_frame, line1(lbl, 11, tc, True), PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
text(s, 10.85, 3.1, 2.0, 0.3, line1("→ 按代分冷热", 11, MUTE))

# --- AFTER bottom: DAMON sampling ---
text(s, 7.05, 3.66, 5.7, 0.32, line1("DAMON · 区域采样", 13, TEALDK, True))
box(s, MSO_SHAPE.RECTANGLE, 7.18, 4.12, 5.42, 0.5, fill=SLATEF, ln=TEAL, lw=1.0)
for dx in (8.26, 9.34, 10.42, 11.5):
    connect(s, dx, 4.12, dx, 4.62, color=TEAL, w=0.75, arrow_end=False)
import itertools
region_centers = [7.72, 8.8, 9.88, 10.96, 12.04]
for rc in region_centers:
    for off in (-0.16, 0.16):
        box(s, MSO_SHAPE.OVAL, rc+off-0.05, 4.32, 0.1, 0.1, fill=TEAL)
text(s, 7.18, 4.74, 5.5, 0.6,
     line1("每区只抽样几页 → 成本与内存大小无关;冷热从「LRU 近似」→「实测访问」", 12, INK))

# --- benefit band ---
box(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, 5.95, 12.33, 1.06, fill=BAND, ln=TEAL, lw=1.0, radius=0.08)
text(s, 0.78, 6.05, 0.9, 0.9, line1("好处", 15, TEALDK, True), anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.7, 6.02, 10.9, 0.95,
     [[{'t':"✓ 让「周期性把内存看一遍」在能耗上可承受 —— 这才「使能」了主动扫描",'s':14,'c':INK}],
      [{'t':"✓ 扫描成本不再随内存规模爆炸　·　kswapd CPU ↓　·　判冷热更准",'s':14,'c':INK}]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=4)

# ============================ SLIDE 3 — innovation 2 =====================
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
text(s, 0.5, 0.3, 12.6, 0.6, line1("创新二 · 触发与控制:被动 → 主动闭环", 30, INK, True))
text(s, 0.5, 0.96, 12.4, 0.4, line1("时机提前 + 触发权逐级放开 → 声明目标·自动调配额的反馈闭环", 14, MUTE))
text(s, 0.55, 1.46, 5.6, 0.4, line1("之前 · 被动撞水位", 16, MUTE, True))
text(s, 6.75, 1.46, 6.0, 0.4, line1("之后 · 主动 + 声明式闭环", 16, TEAL, True))

# --- BEFORE: memory vs time curve ---
connect(s, 1.05, 5.0, 1.05, 2.1, color=SLATE, w=1.5, arrow_end=True)   # y-axis up
connect(s, 1.05, 5.0, 5.95, 5.0, color=SLATE, w=1.5, arrow_end=True)   # x-axis right
text(s, 0.5, 1.92, 1.2, 0.3, line1("内存", 11, MUTE))
text(s, 5.25, 5.05, 1.2, 0.3, line1("时间", 11, MUTE))
connect(s, 1.05, 2.5, 5.85, 2.5, color=RED, w=1.25, dash='dash', arrow_end=False)  # ceiling
text(s, 3.7, 2.2, 2.2, 0.3, line1("内存上限 / 水位", 11, RED))
curve = [(1.05,4.72),(2.1,4.35),(3.0,3.8),(3.7,3.05),(4.2,2.52),(4.85,2.52),(5.4,3.5)]
for (ax,ay),(bx,by) in zip(curve, curve[1:]):
    connect(s, ax, ay, bx, by, color=INK, w=2.0, arrow_end=False)
box(s, MSO_SHAPE.OVAL, 4.42, 2.34, 0.16, 0.16, fill=RED)
text(s, 4.0, 1.95, 2.4, 0.3, line1("✗ 撞顶", 12, RED, True))
text(s, 0.6, 5.28, 5.45, 0.8,
     [[{'t':"✗ ",'s':12.5,'c':RED,'b':True},
       {'t':"压力到了才扫 → 同步 direct reclaim 卡顿 / lmkd 杀;内核被动,外部难干预节奏",'s':12.5,'c':RED}]])

# --- AFTER: timing note + control loop ---
text(s, 6.75, 1.9, 6.3, 0.4, line1("时机:压力『前』周期扫 → 预腾 headroom + 产出冷热信息", 12.5, INK))
nodes = {
 'n1': (6.78, 2.5,  "① 声明固定目标", "(PSI / 内存占用比)", BAND),
 'n2': (10.22,2.5,  "② DAMOS 比例回路调配额", "(硬上限封顶)", BAND),
 'n3': (10.22,4.35, "③ 回收 / 迁移动作", "", WHITE),
 'n4': (6.78, 4.35, "④ 测量当前值", "per-node / per-memcg", WHITE),
}
NW, NH = 2.55, 0.92
for key,(nx,ny,t1,t2,fl) in nodes.items():
    nb = box(s, MSO_SHAPE.ROUNDED_RECTANGLE, nx, ny, NW, NH, fill=fl, ln=TEAL, lw=1.4, radius=0.1)
    paras = [[{'t':t1,'s':12.5,'c':INK,'b':True}]]
    if t2:
        paras.append([{'t':t2,'s':11,'c':MUTE,'f':(MONO if '/' in t2 or 'PSI' in t2 else YAHEI)}])
    fill_tf(nb.text_frame, paras, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, space_after=2)
# clockwise arrows
connect(s, 9.33, 2.96, 10.22, 2.96, color=TEAL, w=1.75, arrow_end=True)   # n1->n2
connect(s, 11.495,3.42, 11.495,4.35, color=TEAL, w=1.75, arrow_end=True)  # n2->n3
connect(s, 10.22,4.81, 9.33, 4.81, color=TEAL, w=1.75, arrow_end=True)    # n3->n4
connect(s, 8.055,4.35, 8.055,3.42, color=TEAL, w=1.75, arrow_end=True)    # n4->n1
text(s, 8.95, 3.62, 1.6, 0.4, line1("反馈闭环", 12, MUTE, True), PP_ALIGN.CENTER)
text(s, 6.75, 5.46, 6.4, 0.4, line1("分工:选谁动手 = 冷热图　·　动多狠 = 目标指标", 12.5, MUTE))

# --- benefit band ---
box(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, 5.9, 12.33, 1.16, fill=BAND, ln=TEAL, lw=1.0, radius=0.07)
text(s, 0.78, 5.98, 0.9, 1.0, line1("好处", 15, TEALDK, True), anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.7, 5.97, 10.9, 1.02,
     [[{'t':"✓ 突发撞进 headroom → 不触顶 → 尾延迟 ↓　·　lmkd 被杀率 ↓",'s':14,'c':INK}],
      [{'t':"（注:削的是「峰值时的阻塞与被杀」,不是峰值内存需求）",'s':11.5,'c':MUTE}],
      [{'t':"✓ 激进度被闭环钳住、免人肉调参　·　硬上限守住续航红线",'s':14,'c':INK}]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=3)

out = "主动扫描两大创新.pptx"
prs.save(out)
print("saved", out)
